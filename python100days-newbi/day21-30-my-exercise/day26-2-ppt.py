"""
pip install python-pptx
"""

from pptx import Presentation 

# create the slide
pres  = Presentation()
# title page
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)

# get title and subtitle
title = slide.shapes.title
subtitle = slide.placeholders[1]
# set the title and subtitle
title.text = "Welcome to my python lesson"
subtitle.text = "use python to creat ppt"
# add slide
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)
# get the shapes in the page
shapes = slide.shapes
# get the title and body
title_shape = shapes.title
body_shape = shapes.placeholders[1]
# edit title
title_shape.text = 'Introduction'
# edit body
# 1.get a text_frame
tf = body_shape.text_frame
tf.text = 'History of Python'
# paragraph level1
p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 1
# paragraph level2
p = tf.add_paragraph()
p.text =  'Guido began to write interpreter for Python.'
p.level = 2

# save 
pres.save("python.pptx")
